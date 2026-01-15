import os.path
import io
import re
from urllib.parse import urlparse, parse_qs
from datetime import datetime
import dateutil.parser as dparser

# Lightweight libraries
from pptx import Presentation
from pdfminer.high_level import extract_text
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Django imports
from django.core.management.base import BaseCommand
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

# Local models
from property.models import PropertyRecord

SCOPES = ['https://www.googleapis.com/auth/drive.readonly']


class Command(BaseCommand):
    help = "Syncs properties; captures dates from parent folders and handles slashes in names."

    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.folder_cache = {}  # To avoid hitting Drive API for the same parent multiple times

    def handle(self, *args, **options):
        # 1. Google Drive Authentication
        creds = None
        if os.path.exists('token.json'):
            creds = Credentials.from_authorized_user_file('token.json', SCOPES)
        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                creds.refresh(Request())
            else:
                flow = InstalledAppFlow.from_client_secrets_file('bdstorage_credentials.json', SCOPES)
                creds = flow.run_local_server(port=0)
            with open('token.json', 'w') as token:
                token.write(creds.to_json())

        service = build('drive', 'v3', credentials=creds)

        # 2. Database Cleanup
        self.stdout.write(self.style.WARNING("Wiping PropertyRecord table for fresh sync..."))
        PropertyRecord.objects.all().delete()

        # 3. Fetch Items
        q = (
            "trashed = false and "
            "createdTime > '2025-01-01T00:00:00Z' and ("
            "mimeType = 'application/vnd.google-apps.folder' or "
            "mimeType contains 'presentation' or "
            "name = 'ai_summary.pdf' or "
            "name = 'recording.mp4'"
            ")"
        )
        all_items = self.get_all_files(service, q)

        # 4. Folder Organization
        folder_data = {}
        for item in all_items:
            if item['mimeType'] == 'application/vnd.google-apps.folder':
                folder_data[item['id']] = {
                    'name': item['name'],
                    'id': item['id'],
                    'parents': item.get('parents', []),
                    'ppt_id': None, 'ppt_link': None,
                    'pdf_id': None, 'pdf_link': None,
                    'mp4_link': None, 'thumb_link': None
                }

        for item in all_items:
            parents = item.get('parents')
            if not parents or parents[0] not in folder_data: continue
            p_id = parents[0]
            name = item.get('name', '').lower()

            if 'presentation' in item.get('mimeType', ''):
                folder_data[p_id]['ppt_id'] = item['id']
                folder_data[p_id]['ppt_link'] = item.get('webViewLink')
                t_link = item.get('thumbnailLink')
                if t_link:
                    folder_data[p_id]['thumb_link'] = t_link.replace('=s220', '=s1000')
            elif name == 'ai_summary.pdf':
                folder_data[p_id]['pdf_id'] = item['id']
                folder_data[p_id]['pdf_link'] = item.get('webViewLink')
            elif name == 'recording.mp4':
                folder_data[p_id]['mp4_link'] = item.get('webViewLink')

        # 5. Extraction and Save
        if not os.path.exists('downloads'): os.makedirs('downloads')
        count = 0

        for f_id, data in folder_data.items():
            if data['ppt_id'] and data['pdf_id']:
                # --- FIX: RETAIN ACTUAL FILENAME WITH SLASHES ---
                # We replace the illegal keyboard '/' with a visual '∕' (Unicode U+2215)
                visual_slash_name = data['name'].replace('/', '∕')

                local_pptx = os.path.join('downloads', f"{visual_slash_name}.pptx")
                local_pdf = os.path.join('downloads', f"{visual_slash_name}_sum.pdf")

                try:
                    self.stdout.write(f"Syncing: {data['name']}")

                    # Search for date in hierarchy
                    presentation_date = self.find_date_in_parents(service, f_id)

                    self.download_file(service, data['ppt_id'], local_pptx)
                    self.download_file(service, data['pdf_id'], local_pdf)

                    retail_url = self.extract_retail_link(local_pptx)
                    prop_id = self.get_property_id(retail_url)
                    extracted_status = self.extract_status_from_pdf(local_pdf)

                    try:
                        ppt_info = self.extract_all_ppt_info(local_pptx)
                    except Exception:
                        ppt_info = {}

                    PropertyRecord.objects.create(
                        property_id=prop_id,
                        presentation_date=presentation_date,  # POPULATED FROM FOLDER HIERARCHY
                        circle=ppt_info.get('circle'),
                        hub=ppt_info.get('hub'),
                        hub_rank=ppt_info.get('hub_rank'),
                        city=ppt_info.get('city'),
                        city_rank=ppt_info.get('city_rank'),
                        final_market_name=ppt_info.get('final_market_name') or data['name'],
                        zone_name=ppt_info.get('zone_name'),
                        ppt_link=data.get('ppt_link'),
                        ai_summary_link=data.get('pdf_link'),
                        recording_link=data.get('mp4_link'),
                        first_slide_image_url=data.get('thumb_link'),
                        status=extracted_status,
                        projected_revenue_lakhs=ppt_info.get('revenue', 'N/A'),
                        total_rent_maintenance=ppt_info.get('rent', 'N/A')
                    )

                    self.stdout.write(self.style.SUCCESS(f"  -> Saved {data['name']} (Date: {presentation_date})"))
                    count += 1

                except Exception as e:
                    self.stdout.write(self.style.ERROR(f"  -> Error on {data['name']}: {e}"))

    def find_date_in_parents(self, service, folder_id):
        """Recursively walks up Drive parents to find a folder name that looks like a date."""
        current_id = folder_id

        while current_id:
            # Check cache first to save API quota
            if current_id in self.folder_cache:
                folder_meta = self.folder_cache[current_id]
            else:
                folder_meta = service.files().get(fileId=current_id, fields="name, parents").execute()
                self.folder_cache[current_id] = folder_meta

            name = folder_meta.get('name', '')

            # Try to parse date from folder name
            try:
                # fuzzy=True handles strings like "Meeting 25 Jan 2015"
                parsed_date = dparser.parse(name, fuzzy=True).date()
                # Sanity check: ensure it's a plausible year for Lenskart properties
                if 2010 < parsed_date.year < 2030:
                    return parsed_date
            except (ValueError, OverflowError):
                pass

            # Move to next parent
            parents = folder_meta.get('parents')
            current_id = parents[0] if parents else None

        return None

    def extract_all_ppt_info(self, path):
        results = {'circle': None, 'hub': None, 'hub_rank': None, 'city': None,
                   'city_rank': None, 'final_market_name': None, 'zone_name': None,
                   'revenue': "N/A", 'rent': "N/A"}
        try:
            prs = Presentation(path)
            if not prs.slides: return results
            first_slide = prs.slides[0]
            slide1_text = "".join([shape.text + "\n" for shape in first_slide.shapes if hasattr(shape, "text")])

            zone_match = re.search(r"ZONE\s*:\s*(.*?)(?:\s*STATE|\s*CITY|\s*PIN CODE|$)", slide1_text,
                                   re.IGNORECASE | re.DOTALL)
            if zone_match:
                results['zone_name'] = re.sub(r'\s*\[Image \d+\]\s*', '', zone_match.group(1)).strip()

            market_match = re.search(r".*?_.*?\(.*?\)_.*?\(.*?\)_.*", slide1_text, re.IGNORECASE)
            if market_match:
                parts = market_match.group(0).strip().split('_')
                if len(parts) >= 4:
                    results['final_market_name'] = parts[-1].strip()
                    for key, part_idx, rank_key in [('city', -2, 'city_rank'), ('hub', -3, 'hub_rank')]:
                        p = parts[part_idx].strip()
                        m = re.search(r"(.*?)\s*\((.*?)\)", p)
                        results[key] = m.group(1).strip() if m else p
                        results[rank_key] = m.group(2).strip() if m else None
                    results['circle'] = re.sub(r'^(Add|BD|Presentation|PPT)\s*', '', parts[-4].strip(),
                                               flags=re.IGNORECASE).strip()

            for slide in prs.slides:
                txt = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"): txt += shape.text + " "
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        for row in shape.table.rows:
                            for cell in row.cells: txt += cell.text + " "
                if results['revenue'] == "N/A":
                    rev_m = re.search(r'GeoIQ Revenue Projection 2025.*?\n?([\d,.]+)', txt, re.IGNORECASE | re.DOTALL)
                    if rev_m: results['revenue'] = rev_m.group(1).replace(',', '')
                if results['rent'] == "N/A":
                    rent_m = re.search(r'Total Rent \+ Maintenance.*?\n?([\d,.]+)', txt, re.IGNORECASE | re.DOTALL)
                    if rent_m: results['rent'] = rent_m.group(1).replace(',', '')
            return results
        except Exception:
            return results

    # ... (extract_status_from_pdf, extract_retail_link, get_property_id, get_all_files, download_file remain the same) ...

    def extract_status_from_pdf(self, path):
        try:
            text = extract_text(path)
            lines = [l.strip() for l in text.split('\n') if l.strip()]
            if not lines: return 'pending'
            context = " ".join(lines[-2:]).lower()
            if 'conditionally approved' in context:
                return 'Conditionally Approved'
            elif 'approved' in context:
                return 'Approved'
            elif any(x in context for x in ['rejected', 'dropped', 'not feasible']):
                return 'Dropped/Rejected'
            elif 'hold' in context:
                return 'Hold'
            return 'pending'
        except:
            return 'pending'

    def extract_retail_link(self, path):
        try:
            prs = Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame: continue
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "RetailIQ" in run.text and run.hyperlink and run.hyperlink.address:
                                return run.hyperlink.address
        except:
            return None
        return None

    def get_property_id(self, url):
        if not url: return None
        try:
            params = parse_qs(urlparse(url).query)
            return params.get('property_id', [None])[0]
        except:
            return None

    def get_all_files(self, service, q):
        items, page_token = [], None
        while True:
            res = service.files().list(q=q,
                                       fields="nextPageToken, files(id, name, webViewLink, mimeType, parents, thumbnailLink)",
                                       pageToken=page_token, pageSize=1000).execute()
            items.extend(res.get('files', []))
            page_token = res.get('nextPageToken')
            if not page_token: break
        return items

    def download_file(self, service, file_id, destination):
        meta = service.files().get(fileId=file_id, fields='mimeType').execute()
        if meta['mimeType'] == 'application/vnd.google-apps.presentation':
            request = service.files().export_media(fileId=file_id,
                                                   mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        else:
            request = service.files().get_media(fileId=file_id)
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done: _, done = downloader.next_chunk()
        with open(destination, 'wb') as f:
            f.write(fh.getvalue())