import os
import re
import io
import django
from pptx import Presentation
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from decimal import Decimal

# --- DJANGO SETUP ---
os.environ.setdefault('DJANGO_SETTINGS_MODULE', 'property_approval_dashboard.settings')
django.setup()

from property.models import PropertyRecord

SCOPES = ['https://www.googleapis.com/auth/drive.readonly']
CLIENT_SECRETS_FILE = 'bdstorage_credentials.json'


def get_drive_service():
    creds = None
    if os.path.exists('token.json'):
        creds = Credentials.from_authorized_user_file('token.json', SCOPES)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(CLIENT_SECRETS_FILE, SCOPES)
            creds = flow.run_local_server(port=0)
        with open('token.json', 'w') as token:
            token.write(creds.to_json())
    return build('drive', 'v3', credentials=creds)


def extract_file_id(url):
    match = re.search(r'/d/([a-zA-Z0-9_-]+)', url)
    return match.group(1) if match else None


def download_and_extract(service, file_id):
    try:
        meta = service.files().get(fileId=file_id).execute()
        mime_type = meta.get('mimeType')
        fh = io.BytesIO()

        if mime_type == 'application/vnd.google-apps.presentation':
            request = service.files().export_media(fileId=file_id,
                                                   mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        else:
            request = service.files().get_media(fileId=file_id)

        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            _, done = downloader.next_chunk()

        fh.seek(0)
        prs = Presentation(fh)
        text_content = ""
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text"):
                text_content += f" {shape.text}"

        lat_raw = re.search(r"LAT:\s*([\d.-]+)", text_content, re.I)
        lon_raw = re.search(r"LON:\s*([\d.-]+)", text_content, re.I)

        return (lat_raw.group(1) if lat_raw else None), (lon_raw.group(1) if lon_raw else None)
    except Exception as e:
        print(f"Error processing {file_id}: {e}")
        return None, None


def run_update():
    service = get_drive_service()
    # Process records where lat/lon is currently empty
    records = PropertyRecord.objects.filter(latitude__isnull=True)

    print(f"Found {records.count()} records to update.")

    for record in records:
        file_id = extract_file_id(record.ppt_link)
        if not file_id:
            continue

        lat_val, lon_val = download_and_extract(service, file_id)

        if lat_val and lon_val:
            try:
                record.latitude = Decimal(lat_val)
                record.longitude = Decimal(lon_val)
                record.save()
                print(f"✅ Updated ID {record.id}: {lat_val}, {lon_val}")
            except Exception as e:
                print(f"⚠️ Could not save ID {record.id}: {e}")
        else:
            print(f"❓ No coordinates found for ID {record.id}")


if __name__ == "__main__":
    run_update()