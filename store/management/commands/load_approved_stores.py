import re
import io
import requests
from django.core.management.base import BaseCommand
from django.core.files.base import ContentFile
from django.utils.dateparse import parse_datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from store.models import ApprovedLenskartStore


class Command(BaseCommand):
    help = 'API Fetch + PPT Download + Contextual Extraction (Using Reference Logic)'

    def get_slide_text(self, slide):
        """Extracts text from shapes and tables within a slide."""
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
        return text

    def handle(self, *args, **options):
        url = "https://retailapis-in.geoiq.ai/bdapp/prod/v1/bd/getAllKissflowProperties"
        headers = {
            'x-api-key': 'eyJ0eXAiOiJKV1QiLCJhbGciOiJIUzI1NiJ9.eyJtYWlsSWRlbnRpdHkiOiJsZW5za2FydC1hZG1pbkBnZW9pcS5pbyJ9.SbaX7dGyaRz_EVwgdE-cQe7_c0pCnM2OPJD2legV0jU'
        }

        self.stdout.write("Fetching properties...")
        try:
            response = requests.post(url, headers=headers, json={}, timeout=60)
            api_data = response.json().get('data', [])
        except Exception as e:
            self.stderr.write(f"API Failed: {e}")
            return

        for item in api_data:
            prop_id = item.get('property_id')

            # 1. Save all API fields to the model
            store, _ = ApprovedLenskartStore.objects.update_or_create(
                property_id=prop_id,
                defaults={
                    'property_name': item.get('property_name'),
                    'property_address': item.get('property_address'),
                    'latitude': item.get('property_lat'),
                    'longitude': item.get('property_lng'),
                    'market_id': item.get('market_id'),
                    'market_name': item.get('market_name'),
                    'kissflow_status': item.get('kissflow_status'),
                    'latest_approval_date': parse_datetime(item.get('latest_approval_date')) if item.get(
                        'latest_approval_date') else None,
                    'property_created_by_info': item.get('property_created_by_info'),
                    'survey_filled_by_info': item.get('survey_filled_by_info'),
                    'ppt_url_raw': item.get('ppt_url'),
                }
            )

            if item.get('ppt_url'):
                try:
                    ppt_res = requests.get(item.get('ppt_url'), timeout=30)
                    if ppt_res.status_code == 200:
                        file_content = ppt_res.content
                        store.ppt_file.save(f"store_{prop_id}.pptx", ContentFile(file_content), save=False)

                        prs = Presentation(io.BytesIO(file_content))

                        # --- 2. EXTRACT SPECS (SLIDE 1) ---
                        if len(prs.slides) > 0:
                            slide1_text = self.get_slide_text(prs.slides[0])

                            # Using the 'Lookahead' logic from your reference script
                            def lookahead_extract(label, next_label, text):
                                pattern = rf'{label}\s*[:\-]?\s*(.*?)(?={next_label})'
                                match = re.search(pattern, text, re.IGNORECASE | re.DOTALL)
                                return match.group(1).strip() if match else "N/A"

                            store.store_size = lookahead_extract("STORE SIZE", "FRONTAGE", slide1_text)
                            store.frontage = lookahead_extract("FRONTAGE", "SIGNAGE WIDTH", slide1_text)
                            store.signage_width = lookahead_extract("SIGNAGE WIDTH", "SIGNAGE HEIGHT", slide1_text)
                            store.signage_height = lookahead_extract("SIGNAGE HEIGHT", "CEILING HEIGHT", slide1_text)
                            store.ceiling_height = lookahead_extract("CEILING HEIGHT", "HEIGHT FROM BEAM BOTTOM",
                                                                     slide1_text)
                            store.height_from_beam_bottom = lookahead_extract("HEIGHT FROM BEAM BOTTOM", "TRADE AREA",
                                                                              slide1_text)
                            store.trade_area = lookahead_extract("TRADE AREA", "Catchment Overview", slide1_text)

                        # --- 3. EXTRACT FINANCIALS (ALL SLIDES) ---
                        all_text = ""
                        for slide in prs.slides:
                            all_text += self.get_slide_text(slide)

                        # Logic for Revenue
                        rev_match = re.search(r'GeoIQ Revenue Projection 2025.*?\n?([\d,.]+)', all_text,
                                              re.IGNORECASE | re.DOTALL)
                        if rev_match:
                            store.geoiq_projected_revenue = rev_match.group(1).replace(',', '')

                        # Logic for Rent
                        rent_match = re.search(r'Total Rent \+ Maintenance.*?\n?([\d,.]+)', all_text,
                                               re.IGNORECASE | re.DOTALL)
                        if rent_match:
                            store.proposed_rent = rent_match.group(1).replace(',', '')

                        store.save()
                        self.stdout.write(f"Updated Store: {prop_id}")

                except Exception as e:
                    self.stderr.write(f"Error on {prop_id}: {e}")

        self.stdout.write(self.style.SUCCESS("Process Complete."))